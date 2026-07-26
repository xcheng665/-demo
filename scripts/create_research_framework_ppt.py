from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE

OUT = '热带海岛建筑环境与智能建造_A3横版可编辑.pptx'
W, H = 16.535, 11.693
INK, GRAY, SOFT, PINK, PALE = '1F1F1F', 'B6B6B6', 'F7F7F7', 'F39AAA', 'FFF3F5'
FONT = 'Microsoft YaHei'

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def rgb(hex_color):
    return RGBColor.from_string(hex_color)

def set_fill(shape, color, trans=0):
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(color); shape.fill.transparency = trans

def set_line(shape, color=GRAY, width=0.65, dash=None):
    shape.line.color.rgb = rgb(color); shape.line.width = Pt(width)
    if dash: shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

def tb(text, x, y, w, h, size=7.2, bold=False, align='center', color=INK, face=FONT, valign='mid', rotate=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background(); box.line.fill.background()
    if rotate is not None: box.rotation = rotate
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if valign == 'mid' else MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.alignment = {'center': PP_ALIGN.CENTER, 'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT}[align]
    run = p.add_run(); run.text = text
    f = run.font; f.name = face; f.size = Pt(size); f.bold = bold; f.color.rgb = rgb(color)
    return box

def rounded(x, y, w, h, fill='FFFFFF', line=GRAY, width=.65, dash=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shp, fill); set_line(shp, line, width, dash); return shp

def rectangle(x, y, w, h, fill='FFFFFF', line=GRAY, width=.65):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shp, fill); set_line(shp, line, width); return shp

def line(x1, y1, x2, y2, color=GRAY, width=.65, dash=False):
    shp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    set_line(shp, color, width, dash); return shp

def arrow(x1, y1, x2, y2, color='4B4B4B'):
    line(x1, y1, x2, y2, color, .7)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x2-.07), Inches(y2-.07), Inches(.14), Inches(.14))
    set_fill(tri, color); set_line(tri, color, .1)
    if abs(x2-x1) >= abs(y2-y1): tri.rotation = 90 if x2 > x1 else 270
    else: tri.rotation = 180 if y2 > y1 else 0

def section(label, x, y, w):
    line(x, y+.12, x+1.16, y+.12, PINK, .85)
    line(x+w-1.16, y+.12, x+w, y+.12, PINK, .85)
    tb(label, x+1.2, y, w-2.4, .22, 8.3, True)

def small(label, x, y, w, h, fs=6.75):
    rounded(x, y, w, h, SOFT, '848484', .55); tb(label, x+.04, y+.02, w-.08, h-.04, fs)

def pinkbox(label, x, y, w, h, fs=7.2):
    rounded(x, y, w, h, PALE, PINK, .7); tb(label, x+.07, y+.03, w-.14, h-.06, fs, True)

def panel(x, y, w, h, zh, en):
    rounded(x, y, w, h, 'FFFFFF', '929292', .78)
    tb(zh, x+.12, y+.13, w-.24, .25, 10.7, True)
    tb(en, x+.12, y+.40, w-.24, .17, 5.75)
    line(x+.1, y+.66, x+w-.1, y+.66, PINK, .75)

def dotted(x, y, w, h):
    rounded(x, y, w, h, 'FFFFFF', '8D8D8D', .6, True)

# Header
tb('热带海岛建筑环境与智能建造', 3.55, .12, 9.45, .47, 23, True)
tb('个人科研实践框架  /  Personal Research Practice Framework', 3.2, .66, 10.15, .22, 10.4)
tb('聚焦热带海岛建筑的可持续性与韧性提升，从低碳材料、建筑知识建模到灾后智能决策，形成“材料—信息—管理—实践”的递进式科研路径。', 2.25, 1.0, 12.05, .18, 6.35)
pinkbox('个人科研实践', 7.16, 1.33, 2.22, .34, 10.3)

py, ph, pw = 1.76, 7.48, 3.84
xs = [.18, 4.27, 8.36, 12.45]
panel(xs[0], py, pw, ph, '研究背景与问题定位', 'Background & Research Focus')
panel(xs[1], py, pw, ph, '材料层：绿色材料与建筑性能', 'Materials Layer')
panel(xs[2], py, pw, ph, '信息层：建筑知识建模与智能检索', 'Information Layer')
panel(xs[3], py, pw, ph, '管理层：灾后恢复与多智能体决策', 'Management Layer')

# Panel 1
section('现实问题', xs[0]+.1, 2.48, pw-.2)
ix, iy, iw, ih, gap = xs[0]+.16, 2.91, .62, 1.36, .13
for i, (ico, label) in enumerate([('♨','高\n温\n高\n湿'),('◔','台\n风\n频\n发'),('♧','资\n源\n约\n束'),('▥','能\n耗\n偏\n高'),('⌘','信\n息\n分\n散')]):
    x = ix+i*(iw+gap); rounded(x,iy,iw,ih,'FFFFFF','A3A3A3',.55)
    tb(ico,x,iy+.1,iw,.22,14,False,face='Segoe UI Symbol')
    tb(label,x+.14,iy+.41,iw-.28,ih-.53,6.7,True)
line(ix+.25,2.77,ix+.25,2.86); line(ix+.25,2.77,ix+3.45,2.77)
for i in range(5): line(ix+i*(iw+gap)+iw/2,2.77,ix+i*(iw+gap)+iw/2,2.87)
line(ix+.25,4.39,ix+.25,4.49,GRAY,.65,True); line(ix+.25,4.49,ix+3.45,4.49,GRAY,.65,True)
section('核心研究问题',xs[0]+.1,4.58,pw-.2); pinkbox('如何通过材料、信息与智能管理\n提升热带海岛建筑的可持续性与韧性？',xs[0]+.29,4.95,pw-.58,.56,7.1)
section('研究定位',xs[0]+.1,5.72,pw-.2); small('热带海岛建筑',xs[0]+1.17,6.02,1.50,.31)
arrow(xs[0]+1.92,6.33,xs[0]+1.92,6.51)
for i,s in enumerate(['绿色低碳材料','建筑知识数字化','灾后智能管理']): small(s,xs[0]+.12+i*1.2,6.58,1.08,.31,5.9)
line(xs[0]+.63,6.51,xs[0]+3.2,6.51)
for v in [.63,1.92,3.2]: line(xs[0]+v,6.51,xs[0]+v,6.57)
line(xs[0]+.63,6.91,xs[0]+3.2,6.91); line(xs[0]+1.92,6.91,xs[0]+1.92,7.04); small('可持续与韧性提升',xs[0]+.95,7.05,2, .32)
section('研究方法导向',xs[0]+.1,7.56,pw-.2)
for i,(ico,lab) in enumerate([('≋','建筑环境'),('♜','实验数据'),('⌘','知识建模'),('◉','智能决策')]):
    x=xs[0]+.19+i*.89; rounded(x,7.89,.72,.74,'FFFFFF',GRAY,.6); tb(ico,x,8.01,.72,.26,14,False,face='Segoe UI Symbol');tb(lab,x+.06,8.36,.60,.20,6.0)

# Panel 2
pinkbox('省级大创项目：椰壳废弃物应用于建筑\n围护结构对建筑能耗的影响研究',xs[1]+.14,2.37,pw-.28,.52,7.1)
section('研究流程',xs[1]+.1,3.05,pw-.2)
for i,s in enumerate(['椰壳废弃物收集','纤维处理与配比设计','试块制备与养护','力学性能测试','围护结构热工分析','建筑能耗影响评估']):
    small(s,xs[1]+.77,3.30+i*.38,2.28,.25,6.5)
    if i<5: arrow(xs[1]+1.91,3.55+i*.38,xs[1]+1.91,3.64+i*.38)
section('关键内容',xs[1]+.1,5.61,pw-.2)
kx=[xs[1]+.16,xs[1]+1.49,xs[1]+2.79]
for i,s in enumerate(['材料制备','力学测试','热工/能耗分析']): small(s,kx[i],5.96,1.03,.28,6.1)
for i,s in enumerate(['椰壳纤维\n混凝土配合比\n资源化利用','抗压\n抗折\n剪裂','传热性能\n围护结构\n节能效益']): dotted(kx[i]-.04,6.37,1.11,.93); tb(s,kx[i]+.05,6.49,.93,.65,6.1)
section('阶段成果',xs[1]+.1,7.55,pw-.2)
for i,s in enumerate(['实验数据','性能认识','绿色材料应用基础']):
    small(s,xs[1]+.15+i*1.34,7.89,1.06,.47,6.15)
    if i<2: arrow(xs[1]+1.22+i*1.34,8.12,xs[1]+1.43+i*1.34,8.12)
line(xs[1]+.1,8.82,xs[1]+pw-.1,8.82,PINK,.75);tb('关键词：建筑物理 / 低碳材料 / 热带建筑',xs[1]+.15,8.92,pw-.3,.16,6.15)

# Panel 3
pinkbox('建筑规范可视化检索与校验系统',xs[2]+.14,2.39,pw-.28,.37,7.3)
section('技术流程',xs[2]+.1,2.96,pw-.2)
for i,s in enumerate(['规范文本收集','条文切片与结构化','实体—属性—关系提取','知识网络构建','交互式可视化','LLM 辅助检索与校验']):
    small(s,xs[2]+.93,3.20+i*.38,1.98,.25,6.25)
    if i<5: arrow(xs[2]+1.92,3.45+i*.38,xs[2]+1.92,3.54+i*.38)
section('系统组成',xs[2]+.1,5.52,pw-.2)
for i,(title,items) in enumerate([('数据层',['规范文本','指标条文','构件信息']),('建模层',['Python处理','NetworkX','Pyvis']),('应用层',['语义检索','规范校验','知识追溯'])]):
    x=xs[2]+.13+i*1.2; dotted(x,5.86,1.07,1.61);tb(title,x+.03,6.00,1.01,.18,7,True)
    for j,s in enumerate(items):small(s,x+.12,6.28+j*.34,.83,.25,5.45)
section('信息层作用',xs[2]+.1,7.54,pw-.2)
for i,(ico,lab) in enumerate([('▤','材料性能数据'),('⌘','知识结构化'),('◯','可查询/\n可推理'),('▣','支撑设计\n与管理')]):
    x=xs[2]+.23+i*.91;tb(ico,x,7.89,.45,.28,16,False,face='Segoe UI Symbol');tb(lab,x-.17,8.27,.80,.34,5.5)
    if i<3:arrow(x+.48,8.06,x+.74,8.06)
line(xs[2]+.1,8.82,xs[2]+pw-.1,8.82,PINK,.75);tb('关键词：知识图谱 / 可视化 / LLM',xs[2]+.15,8.92,pw-.3,.16,6.15)

# Panel 4
pinkbox('台风灾后建筑损伤与恢复多智能体研究',xs[3]+.12,2.40,pw-.24,.37,7.15)
section('决策流程',xs[3]+.1,2.98,pw-.2)
for i,s in enumerate(['灾损多源数据输入','损伤信息识别','任务与约束建模','多智能体协同','恢复计划生成']):
    small(s,xs[3]+1.00,3.23+i*.38,1.80,.25,6.25)
    if i<4:arrow(xs[3]+1.90,3.48+i*.38,xs[3]+1.90,3.57+i*.38)
section('多智能体系统',xs[3]+.1,5.28,pw-.2);small('恢复决策系统',xs[3]+1.19,5.56,1.45,.29)
line(xs[3]+1.92,5.85,xs[3]+1.92,6.03);line(xs[3]+.45,6.03,xs[3]+3.39,6.03)
for i,(a,b) in enumerate([('信息整合\nAgent','修缮优先级'),('约束识别\nAgent','材料推荐'),('任务规划\nAgent','工作计划'),('资源协调\nAgent','资源配置')]):
    x=xs[3]+.10+i*.91;line(x+.42,6.03,x+.42,6.16);pinkbox(a,x,6.18,.80,.44,5.4);arrow(x+.40,6.63,x+.40,6.89);small(b,x,6.91,.80,.37,5.4)
section('管理层价值',xs[3]+.1,7.53,pw-.2)
for i,(ico,lab) in enumerate([('▢','知识支持'),('☷','约束识别'),('♧','协同决策'),('⌁','灾后恢复\n效率提升')]):
    x=xs[3]+.27+i*.87;tb(ico,x,7.86,.32,.28,15,False,face='Segoe UI Symbol');tb(lab,x-.19,8.25,.70,.31,5.5)
    if i<3:arrow(x+.38,8.06,x+.67,8.06)
line(xs[3]+.1,8.82,xs[3]+pw-.1,8.82,PINK,.75);tb('关键词：台风灾后 / 多智能体 / 智能建造',xs[3]+.12,8.92,pw-.24,.16,5.9)

# Bottom
rounded(.18,9.39,W-.36,2.05,'FFFFFF','929292',.75)
tb('融合应用与能力递进',5.25,9.51,6,.26,12,True)
for i,(ico,lab) in enumerate([('♜','材料实验'),('▤','数据治理'),('⌘','知识建模'),('◉','智能决策'),('▥','设计应用')]):
    x=.78+i*2.47; circ=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(9.84),Inches(.61),Inches(.61));set_fill(circ,'FFFFFF');set_line(circ,'E5B3BE',.75)
    tb(ico,x,9.96,.61,.22,15,False,face='Segoe UI Symbol');tb(lab,x-.17,10.59,.95,.19,7.1,True)
    if i<4:arrow(x+.78,10.15,x+2.18,10.15)
tb('形成面向热带海岛建筑的跨学科科研路径：材料提供物质基础，信息建立知识桥梁，管理支撑智能决策。',.98,11.06,9.82,.18,6.9,False,'left')
rounded(11.44,9.72,4.77,1.38,'FFFFFF',PINK,.8);tb('关键工具与平台',12.58,9.67,2.55,.22,8,True)
line(11.55,9.75,12.52,9.75,PINK,.75);line(15.68,9.75,16.1,9.75,PINK,.75)
for i,s in enumerate(['Py','⌘','⌁','◉','≋','▥','AIGC']):tb(s,11.66+i*.62,10.13,.38,.28,6.5 if i==6 else 14,False,face='Segoe UI Symbol')
tb('Python  /  NetworkX  /  Pyvis  /  LLM  /  建筑物理  /  BIM  /  AIGC',11.64,10.72,4.32,.18,6.3)
tb('A3 横向 · 全部由可编辑 PowerPoint 对象构成',.25,11.48,4.1,.10,3.8,False,'left','999999')
prs.save(OUT)
